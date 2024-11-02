from flask import Flask, render_template, request, send_file
import os
from pdf2docx import Converter
from pptx import Presentation
from pdf2image import convert_from_path
import fitz  # PyMuPDF
from moviepy.editor import VideoFileClip
from PIL import Image
import pandas as pd
from docx import Document  # Import Document class
from fpdf import FPDF

app = Flask(__name__)

# Function to convert PDF to DOCX
def pdf_to_docx(pdf_path, docx_path):
    cv = Converter(pdf_path)
    cv.convert(docx_path, start=0, end=None)
    cv.close()

# Function to convert DOCX to PDF
def docx_to_pdf(docx_path, pdf_path):
    doc = Document(docx_path)
    pdf = fitz.open()
    for para in doc.paragraphs:
        pdf.insert_page(-1, text=para.text)
    pdf.save(pdf_path)

# Function to convert image to PDF
def image_to_pdf(image_path, pdf_path):
    image = Image.open(image_path)
    image.convert("RGB").save(pdf_path, "PDF")

# Function to convert Excel to CSV
def excel_to_csv(excel_path, csv_path):
    df = pd.read_excel(excel_path)
    df.to_csv(csv_path, index=False)

# Function to convert CSV to Excel
def csv_to_excel(csv_path, excel_path):
    df = pd.read_csv(csv_path)
    df.to_excel(excel_path, index=False)

# Function to convert Video to Audio
def video_to_audio(video_path, audio_path):
    clip = VideoFileClip(video_path)
    clip.audio.write_audiofile(audio_path)

# Function to convert PPTX to PDF
def pptx_to_pdf(pptx_path, pdf_path):
    presentation = Presentation(pptx_path)
    pdf = FPDF()
    
    for slide in presentation.slides:
        pdf.add_page()
        
        # Create an image from the slide
        slide_image_path = "temp_slide.jpg"
        slide.shapes[0].element.getparent().remove(slide.shapes[0].element)  # Remove first shape
        slide.shapes.add_picture(slide_image_path, 0, 0, width=pdf.w, height=pdf.h)
        
        # Save slide as an image and add to PDF
        slide.shapes._spTree.remove(slide.shapes._spTree[0])  # This removes the slide background
        pdf.image(slide_image_path, 0, 0, 210, 297)  # A4 size

    pdf.output(pdf_path)

# Function to convert PDF to PPTX
def pdf_to_pptx(pdf_path, pptx_path):
    pdf_document = fitz.open(pdf_path)
    presentation = Presentation()

    for page_num in range(len(pdf_document)):
        slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        # Extract text from the PDF
        text = pdf_document[page_num].get_text()
        
        # Add text to slide
        if text:
            textbox = slide.shapes.add_textbox(left=0, top=0, width=Inches(10), height=Inches(7.5))
            textbox.text = text

        # Convert PDF page to image for adding to slide
        image_list = convert_from_path(pdf_path, first_page=page_num + 1, last_page=page_num + 1)
        for img in image_list:
            img_path = f"temp_image_{page_num}.png"
            img.save(img_path, "PNG")
            slide.shapes.add_picture(img_path, 0, 0, width=Inches(10), height=Inches(7.5))
            os.remove(img_path)  # Remove temp image after adding to slide

    presentation.save(pptx_path)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    uploaded_file = request.files['file']
    output_format = request.form['output_format']
    input_filename = uploaded_file.filename
    input_path = os.path.join("uploads", input_filename)
    uploaded_file.save(input_path)

    base, ext = os.path.splitext(input_filename)
    output_path = os.path.join("outputs", f"{base}.{output_format}")

    # Conversion logic based on selected output format
    if ext == ".pdf" and output_format == "docx":
        pdf_to_docx(input_path, output_path)
    elif ext == ".docx" and output_format == "pdf":
        docx_to_pdf(input_path, output_path)
    elif ext in [".jpg", ".png"] and output_format == "pdf":
        image_to_pdf(input_path, output_path)
    elif ext == ".xlsx" and output_format == "csv":
        excel_to_csv(input_path, output_path)
    elif ext == ".csv" and output_format == "xlsx":
        csv_to_excel(input_path, output_path)
    elif ext in [".mp4", ".mov"] and output_format == "mp3":
        video_to_audio(input_path, output_path)
    elif ext == ".pptx" and output_format == "pdf":
        pptx_to_pdf(input_path, output_path)
    elif ext == ".pdf" and output_format == "pptx":
        pdf_to_pptx(input_path, output_path)
    else:
        return "Unsupported file type or conversion."

    return send_file(output_path, as_attachment=True)

if __name__ == '__main__':
    os.makedirs("uploads", exist_ok=True)
    os.makedirs("outputs", exist_ok=True)
    app.run(host='0.0.0.0', port=8721, debug=True)
